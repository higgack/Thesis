import asyncio
import logging
import re
from pathlib import Path
import httpx
import trafilatura
from pypdf import PdfReader
from youtube_transcript_api import YouTubeTranscriptApi

log = logging.getLogger(__name__)

_YOUTUBE_RE = re.compile(
    r"(?:youtube\.com/watch\?v=|youtu\.be/|youtube\.com/shorts/)([\w-]{11})"
)
_ARXIV_RE = re.compile(r"arxiv\.org/(?:abs|pdf)/([\w.\-]+)")


def is_youtube(url: str) -> str | None:
    m = _YOUTUBE_RE.search(url)
    return m.group(1) if m else None


# Blog platforms whose URLs we tag as source_type "blog" (vs generic
# "web") so the study-notes dashboard can filter them separately. Host
# substring match — covers the common KR/global blog hosts.
_BLOG_HOST_SUBSTR = (
    "blog.naver.com", "tistory.com", "brunch.co.kr", "velog.io",
    "medium.com", "blogspot.com", "blogger.com", "wordpress.com",
    "post.naver.com", "brunchstory.co",
)


def is_blog(url: str) -> bool:
    try:
        from urllib.parse import urlparse
        host = urlparse(url).netloc.lower()
    except Exception:
        return False
    return any(h in host for h in _BLOG_HOST_SUBSTR)


_JS_PLACEHOLDER_PATTERNS = (
    "javascript is not available",
    "please enable javascript",
    "you need to enable javascript",
    "javascript이 비활성화",
)
_MIN_BODY_CHARS = 200


# Hosts that block bot/datacenter access — skip the fetch entirely so
# the user gets the recovery guidance (📋 본문 복사 / 스크린샷) instantly
# instead of waiting 30+ seconds for a TCP timeout. Includes:
#   • social/auth-walled (LinkedIn, Facebook, etc.)
#   • paywalled financial press (Reuters, Bloomberg, WSJ, FT, etc.)
_BLOCKED_HOSTS = (
    # Auth walls — bot can never see the body
    "linkedin.com", "facebook.com", "instagram.com", "story.kakao.com",
    # X (twitter) — public oEmbed strips threads + we don't pay for the
    # API. Forwarded channels routinely cite X posts; without this block
    # every cite spawns a useless ingest that yields an empty body and
    # spends a retry slot.
    "x.com", "twitter.com",
    # English-language paywalls — server returns a stub or login redirect
    # (한국 매체는 보통 본문 추출되니 제외; 막으면 좋은 자료까지 잃음)
    "reuters.com", "bloomberg.com", "ft.com", "wsj.com",
    "economist.com", "nytimes.com", "washingtonpost.com", "barrons.com",
    # Korean URL shorteners — opaque redirect, trafilatura can't extract
    # anything from the landing page. Every digest cite of these was
    # piling up in the failed log without ever succeeding.
    "buly.kr", "vo.la", "zrr.kr", "bit.ly", "tinyurl.com",
    # Forum / community boards with low signal + heavy noise (long
    # comment threads that overwhelm any analyst content).
    "dvdprime.com",
    # Stub / placeholder hosts that extract no actual content:
    #   • finance.naver.com / n.stock.naver.com — body is just the
    #     "증권사 로그인 안내" boilerplate; the real price/financials
    #     are loaded by JS and don't surface to trafilatura.
    #   • dart.fss.or.kr — viewer URL extracts only the report table
    #     of contents; substantive text lives in the attached PDF
    #     which needs to be downloaded separately. Daju channel
    #     already strips these via _PLAIN_URL_STRIP_PATTERNS so
    #     blocking here has no downside.
    "finance.naver.com", "n.stock.naver.com", "stock.naver.com",
    "dart.fss.or.kr",
)


def _is_blocked_host(url: str) -> bool:
    """True if url's host matches the _BLOCKED_HOSTS list. Used by
    pipeline.ingest_url to short-circuit BEFORE the failure log gets
    touched, so blocked URLs no longer accumulate in /failed."""
    try:
        from urllib.parse import urlparse
        host = urlparse(url).netloc.lower()
        return any(h in host for h in _BLOCKED_HOSTS)
    except Exception:
        return False


# Realistic desktop Chrome UA. Naver (and several other sites) serve a
# stub/redirect to a bare "Mozilla/5.0 SecondBrain" UA but the full post
# HTML to a normal browser string. A realistic UA is strictly more
# widely accepted than a custom one, so it doubles as a small recall win
# for other hosts too.
_BROWSER_UA = (
    "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
    "(KHTML, like Gecko) Chrome/124.0.0.0 Safari/537.36"
)

# URL shortener hosts. We resolve these to their final destination
# BEFORE the rest of the ingest pipeline runs so dedup, blocked-host
# check, and body fetch operate on the real URL. Without this, every
# buly.kr-laced Korean digest gets all its URLs silent-blocked (the
# shortener landing page yields no body to trafilatura). After
# unshortening, the real destination flows through normal rules — if
# THAT host is also in _BLOCKED_HOSTS (e.g. t.co → twitter.com) the
# usual block still catches it.
_URL_SHORTENERS = (
    "buly.kr", "vo.la", "zrr.kr", "han.gl",        # KR
    "bit.ly", "tinyurl.com", "is.gd", "ow.ly",     # generic
    "t.co", "lnkd.in", "goo.gl",                   # social/google
)


async def unshorten_url(url: str) -> str:
    """If url's host is a known shortener, follow redirects (HEAD →
    cheap, no body) to the real destination. Returns the original url
    on any failure or non-shortener input. Bounded by a short timeout
    so a slow/dead shortener can't stall ingest."""
    try:
        from urllib.parse import urlparse
        host = urlparse(url).netloc.lower()
    except Exception:
        return url
    if not any(host == s or host.endswith("." + s) for s in _URL_SHORTENERS):
        return url
    try:
        async with httpx.AsyncClient(
            timeout=10, follow_redirects=True,
            headers={"User-Agent": _BROWSER_UA},
        ) as c:
            try:
                r = await c.head(url)
            except Exception:
                r = await c.get(url)
            else:
                # Some shorteners refuse HEAD — fall back to GET.
                if r.status_code in (400, 403, 405):
                    r = await c.get(url)
        final = str(r.url)
        if final and final != url:
            log.info("unshortened %s → %s", url[:80], final[:120])
            return final
    except Exception as e:
        log.debug("unshorten failed for %s: %s", url[:80], e)
    return url


_NAVER_BLOG_RE = re.compile(
    r"^https?://(?:m\.)?blog\.naver\.com/([^/?#]+)/(\d+)(?:[/?#].*)?$"
)


def _rewrite_for_better_fetch(url: str) -> str:
    """Rewrite known JS-iframe / SPA URLs to alternative forms that
    yield real HTML body to trafilatura. Currently:
      • blog.naver.com/<user>/<id>  →
        blog.naver.com/PostView.naver?blogId=<user>&logNo=<id>
        The desktop post page is just a shell whose <iframe id=mainFrame>
        loads this PostView endpoint; the mobile (m.blog) page renders
        the body via JS. PostView ships the full post HTML inline, so
        trafilatura sees real markup. (Was m.blog rewrite — that page
        also lazy-loads the body via JS and extracted empty, which
        auto-blocked the whole host.)
    Returns the original URL when no rewrite rule matches."""
    m = _NAVER_BLOG_RE.match(url)
    if m:
        user, post_id = m.group(1), m.group(2)
        return (f"https://blog.naver.com/PostView.naver?"
                f"blogId={user}&logNo={post_id}")
    return url


async def load_url(url: str) -> tuple[str, str, str | None, list[str]]:
    """Returns (title, text, hint_summary, outlinks).

    hint_summary is a free, source-provided abstract / og:description that
    can replace LLM summarization when good enough.

    outlinks are http(s) URLs the author placed in the article BODY (site
    chrome/nav excluded), surfaced for the optional in-post link prompt.
    Empty for youtube / arxiv / pdf / blocked-host results."""
    yt = is_youtube(url)
    if yt:
        t, b, h = await load_youtube(yt, url)
        return t, b, h, []
    if m := _ARXIV_RE.search(url):
        t, b, h = await load_arxiv(m.group(1))
        return t, b, h, []
    if _is_reddit(url):
        # Reddit's public .json endpoint returns the post + comments without
        # JS/auth — Phase-0 style. Falls through to the normal fetch if it's
        # rate-limited (403), so nothing is lost.
        rt, rb, rh = await _load_reddit(url)
        if rb and len(rb) >= _MIN_BODY_CHARS:
            return rt or url[:200], rb, rh, []
    try:
        from urllib.parse import urlparse
        host = urlparse(url).netloc.lower()
        if any(h in host for h in _BLOCKED_HOSTS):
            return "", "", None, []  # short-circuit; downstream says empty
    except Exception:
        pass
    # Rewrite SPA-heavy URLs to scraper-friendly alternates BEFORE the
    # fetch so trafilatura sees real markup. Naver blog is the big
    # one: desktop ships an iframe, mobile ships inline content.
    fetch_url = _rewrite_for_better_fetch(url)
    headers = {"User-Agent": _BROWSER_UA}
    # Naver PostView serves the body only when the Referer points at the
    # owning post; supply the original URL the user forwarded.
    if fetch_url != url and "blog.naver.com" in fetch_url:
        headers["Referer"] = url
    title, body, hint = "", "", None
    links: list[str] = []
    html = ""
    try:
        async with httpx.AsyncClient(timeout=30, follow_redirects=True,
                                     headers=headers) as c:
            r = await c.get(fetch_url)
            r.raise_for_status()
            ct = r.headers.get("content-type", "").lower()
            if "application/pdf" in ct or r.content[:5] == b"%PDF-":
                t, b, h = await _load_pdf_from_bytes(r.content, str(r.url))
                return t, b, h, []
            html = r.text
        title, body, hint = await asyncio.to_thread(_parse_html, fetch_url, html)
        links = await asyncio.to_thread(
            _extract_content_links, html, str(fetch_url))
    except Exception:
        pass

    # trafilatura routinely misses SmartEditor / many blog & news CMS
    # layouts, so when the body comes back thin, pull the largest article
    # container straight from the DOM before the flaky JS renderer. Helps
    # Naver(.se-main-container), Tistory, brunch, news sites, etc. — and
    # since load_url is shared, this improves BOTH study notes and the
    # main RAG ingest.
    if _is_js_placeholder(body) and html:
        dom = await asyncio.to_thread(_extract_dom_body, html)
        if dom:
            body = dom

    # Still thin → retry the fetch with curl_cffi's real browser TLS
    # fingerprint. Clears Cloudflare/WAF 403s that the UA-only httpx call
    # above can't (and covers the case where httpx raised, html == "").
    if _is_js_placeholder(body):
        cffi_html = await _fetch_html_impersonate(fetch_url, headers)
        if cffi_html:
            c_title, c_body, c_hint = await asyncio.to_thread(
                _parse_html, fetch_url, cffi_html)
            if _is_js_placeholder(c_body):
                dom2 = await asyncio.to_thread(_extract_dom_body, cffi_html)
                if dom2:
                    c_body = dom2
            if c_body and len(c_body) >= _MIN_BODY_CHARS:
                title = c_title or title
                body = c_body
                hint = c_hint or hint
                if not links:
                    links = await asyncio.to_thread(
                        _extract_content_links, cffi_html, str(fetch_url))

    if _is_js_placeholder(body):
        try:
            # Use the rewritten URL (e.g. Naver PostView) so Jina renders
            # the content endpoint, not the JS-only wrapper page.
            j_title, j_body, j_hint = await _load_via_jina(fetch_url)
            if j_body and len(j_body) >= _MIN_BODY_CHARS:
                title = j_title or title
                body = j_body
                hint = j_hint or hint
        except Exception:
            pass

    # Last resort: Jina with its real headless-browser engine. Slower but
    # renders stubborn SPA pages (e.g. notion.site) that the default Jina
    # mode returns empty for.
    if _is_js_placeholder(body):
        try:
            j_title, j_body, j_hint = await _load_via_jina(
                fetch_url, engine="browser")
            if j_body and len(j_body) >= _MIN_BODY_CHARS:
                title = j_title or title
                body = j_body
                hint = j_hint or hint
        except Exception:
            pass

    return title or url[:200], body, hint, links


# Article/body containers across common blog & CMS platforms, in rough
# priority order. Naver SmartEditor first (its wrapper is unmistakable),
# then generic article roots. Used only as a fallback when trafilatura
# returns a thin/placeholder body, and gated on _MIN_BODY_CHARS, so a bad
# match can't override a good trafilatura extraction.
_ARTICLE_SELECTORS = (
    ".se-main-container",            # Naver SmartEditor ONE (current)
    "#postViewArea", "div.post_ct",  # Naver legacy editors
    "div.se_component_wrap",
    "#dic_area", "#articleBodyContents", "#newsct_article",  # Naver News
    "#article-view-content-div", ".article_view", "#articeBody",  # KR news CMS
    "#news_body_area", ".news_end", "#CmAdContent",
    "[itemprop='articleBody']",
    "article",
    ".article-body", ".article_body", ".article-content", ".articleView",
    ".entry-content", ".post-content", ".post-body", ".tt_article_useless_p_margin",
    ".contents_style", ".news-content", ".view-content", ".board-view",
    ".story-body", ".content__body", ".article__body", "main",
)


def _extract_dom_body(html: str) -> str:
    """Fallback body extraction from raw HTML for blog/CMS layouts that
    trafilatura misses (Naver SmartEditor, Tistory skins, news CMS, …).
    Strips obvious chrome, then returns the largest matching content
    container's text (or '' when nothing substantial is found)."""
    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(html, "html.parser")
    except Exception:
        return ""
    for tag in soup(["script", "style", "nav", "header", "footer",
                     "aside", "form", "noscript"]):
        tag.decompose()
    best = ""
    for sel in _ARTICLE_SELECTORS:
        try:
            nodes = soup.select(sel)
        except Exception:
            continue
        for node in nodes:
            txt = node.get_text("\n", strip=True)
            if len(txt) > len(best):
                best = txt
    return best if len(best) >= _MIN_BODY_CHARS else ""


def _is_js_placeholder(body: str) -> bool:
    if len(body) < _MIN_BODY_CHARS:
        return True
    low = body.lower()
    return any(p in low for p in _JS_PLACEHOLDER_PATTERNS)


async def _load_via_jina(url: str,
                         engine: str | None = None) -> tuple[str, str, str | None]:
    """r.jina.ai renders JS and returns clean markdown. Free, no auth.

    engine="browser" forces Jina's full headless-browser engine — slower,
    but renders stubborn SPA pages (Notion etc.) the default mode misses.
    """
    api_url = f"https://r.jina.ai/{url}"
    headers = {"Accept": "text/plain"}
    timeout = 60.0
    if engine:
        headers["X-Engine"] = engine          # "browser" = real headless render
        headers["X-Timeout"] = "30"
        timeout = 90.0
    async with httpx.AsyncClient(timeout=timeout, follow_redirects=True,
                                 headers=headers) as c:
        r = await c.get(api_url)
        r.raise_for_status()
        text = r.text
    title = ""
    body = text
    if "Markdown Content:" in text:
        head, body = text.split("Markdown Content:", 1)
        body = body.strip()
        if "Title:" in head:
            title = head.split("Title:", 1)[1].split("\n", 1)[0].strip()
    return title, body, None


async def _fetch_html_impersonate(url: str, headers: dict) -> str:
    """Re-fetch with curl_cffi forging a real Chrome TLS/JA3 fingerprint.
    Bypasses many Cloudflare/WAF 403s that the UA-only httpx call can't.
    Returns html text, or '' on any failure / non-HTML / blocked status.
    Best-effort: a missing curl_cffi dep just yields ''."""
    def _go() -> str:
        try:
            from curl_cffi import requests as _cffi
        except Exception:
            return ""
        try:
            r = _cffi.get(url, headers=headers, impersonate="chrome",
                          timeout=30, allow_redirects=True)
        except Exception:
            return ""
        if r.status_code >= 400:
            return ""
        ct = (r.headers.get("content-type") or "").lower()
        if ct and "html" not in ct and "text" not in ct:
            return ""  # PDFs etc. are handled by the primary httpx path
        try:
            return r.text or ""
        except Exception:
            return ""
    return await asyncio.to_thread(_go)


_REDDIT_RE = re.compile(r"(?:^|\.)reddit\.com$")


def _is_reddit(url: str) -> bool:
    try:
        from urllib.parse import urlparse
        return bool(_REDDIT_RE.search(urlparse(url).netloc.lower()))
    except Exception:
        return False


async def _load_reddit(url: str) -> tuple[str, str, str | None]:
    """Reddit exposes any post at <url>.json (post + comments, no auth).
    Returns (title, body, None); ('', '', None) on block/parse failure so
    the caller falls back to the normal fetch path."""
    base = url.split("?")[0].split("#")[0].rstrip("/")
    api = base if base.endswith(".json") else base + ".json"
    headers = {"User-Agent": _BROWSER_UA}
    try:
        async with httpx.AsyncClient(timeout=30, follow_redirects=True,
                                     headers=headers) as c:
            r = await c.get(api)
            r.raise_for_status()
            data = r.json()
    except Exception:
        return "", "", None
    try:
        post = data[0]["data"]["children"][0]["data"]
        title = (post.get("title") or "").strip()
        parts: list[str] = []
        if (post.get("selftext") or "").strip():
            parts.append(post["selftext"].strip())
        if len(data) > 1:
            for ch in data[1]["data"]["children"][:25]:
                d = ch.get("data", {})
                cbody = (d.get("body") or "").strip()
                author = d.get("author") or ""
                if cbody and author not in ("AutoModerator", "", "[deleted]"):
                    parts.append(f"— {author}: {cbody}")
        return title, "\n\n".join(parts).strip(), None
    except Exception:
        return "", "", None


def _parse_html(url: str, html: str) -> tuple[str, str, str | None]:
    # Recall-favoring extraction: news/blog pages with mixed markup often
    # get under-extracted by the default precision mode. favor_recall +
    # tables pulls the full article (and financial/data tables) without the
    # comment sections. Fall back to defaults if the installed trafilatura
    # predates these kwargs so the call can never break.
    try:
        extracted = trafilatura.extract(
            html, url=url, include_comments=False, include_tables=True,
            favor_recall=True) or ""
    except TypeError:
        extracted = trafilatura.extract(html, include_comments=False) or ""
    meta = trafilatura.extract_metadata(html)
    title = (meta.title if meta and meta.title else url)[:200]
    hint = (meta.description if meta and meta.description else None)
    return title, extracted.strip(), hint


_MAX_OUTLINKS = 12
_OUTLINK_IMG_EXT = (".jpg", ".jpeg", ".png", ".gif", ".webp", ".svg", ".bmp")
# Asset / CDN / tracker hosts that are never article content — Naver
# image stores, profile thumbnails, share widgets, etc.
_OUTLINK_SKIP_SUBSTR = (
    "pstatic.net", "phinf.naver", "blogpfthumb", "ssl.pstatic",
    "naver.com/profile", "/share", "javascript:",
)


def _extract_content_links(html: str, base_url: str) -> list[dict]:
    """In-post links the author placed in the article BODY (not nav /
    menu / related-widget chrome). Runs trafilatura with include_links
    so the content detector drops site chrome, then pulls absolute URLs
    (and their anchor text) out of the markdown. Images, asset/tracker
    hosts, and the page's own URL are filtered. Returns a list of
    {"url","anchor"} — deduped, order-preserved, capped at _MAX_OUTLINKS.
    anchor is the link's display text (often the destination's title for
    Naver related-post cards), used as a free preview fallback."""
    try:
        md = trafilatura.extract(
            html, include_links=True, include_comments=False,
            output_format="markdown",
        ) or ""
    except Exception:
        return []
    from urllib.parse import urlparse
    base = urlparse(base_url)
    out: list[dict] = []
    seen: set[str] = set()
    # Capture both the anchor text and the URL: [anchor](url)
    for m in re.finditer(r"\[([^\]]*)\]\((https?://[^\s)]+)\)", md):
        anchor = (m.group(1) or "").strip()
        u = m.group(2).rstrip(".,)")
        low = u.lower()
        if low.endswith(_OUTLINK_IMG_EXT):
            continue
        if any(s in low for s in _OUTLINK_SKIP_SUBSTR):
            continue
        if u in seen:
            continue
        p = urlparse(u)
        # Skip the page's own URL (self-link / canonical echo).
        if p.netloc.lower() == base.netloc.lower() and p.path == base.path:
            continue
        seen.add(u)
        out.append({"url": u, "anchor": anchor[:120]})
        if len(out) >= _MAX_OUTLINKS:
            break
    return out


def _parse_preview_meta(html: str) -> dict:
    """Title + description from a page's metadata only (og: tags /
    <title> / meta description) — no body extraction. Cheap link-card
    style preview."""
    try:
        meta = trafilatura.extract_metadata(html)
    except Exception:
        meta = None
    title = (meta.title if meta and meta.title else "") or ""
    desc = (meta.description if meta and meta.description else "") or ""
    return {"title": title.strip()[:120], "desc": desc.strip()[:200]}


async def fetch_link_preview(url: str) -> dict:
    """Lightweight {"title","desc"} preview for a link WITHOUT learning
    it: one short GET + metadata parse. No Jina, no body extraction, no
    cost. Short timeout so a slow host can't stall the prompt; failures
    return empties and the caller falls back to the anchor text."""
    fetch_url = _rewrite_for_better_fetch(url)
    headers = {"User-Agent": _BROWSER_UA}
    if fetch_url != url and "blog.naver.com" in fetch_url:
        headers["Referer"] = url
    try:
        async with httpx.AsyncClient(timeout=8, follow_redirects=True,
                                     headers=headers) as c:
            r = await c.get(fetch_url)
            r.raise_for_status()
            ct = r.headers.get("content-type", "").lower()
            if "text/html" not in ct and "application/xhtml" not in ct:
                return {"title": "", "desc": ""}
            html = r.text[:300_000]
    except Exception:
        return {"title": "", "desc": ""}
    return await asyncio.to_thread(_parse_preview_meta, html)


async def _fetch_youtube_subs_yt_dlp(video_id: str) -> str:
    """Fallback transcript fetcher using yt-dlp. yt-dlp talks to
    YouTube via the same internal endpoints the web client uses, so
    it usually works from GCP/AWS server IPs where the
    youtube-transcript-api library gets blocked. Tries manual
    Korean → English subtitles first, then auto-generated captions
    in the same order. Returns the plain transcript text (one line
    per caption block) or '' if nothing is available."""
    try:
        from yt_dlp import YoutubeDL
    except Exception:
        log.warning("yt-dlp not installed — skipping fallback")
        return ""

    import json
    url = f"https://www.youtube.com/watch?v={video_id}"
    opts = {
        "skip_download": True,
        "quiet": True,
        "no_warnings": True,
        "writesubtitles": False,
        "writeautomaticsub": False,
        "subtitleslangs": ["ko", "en"],
        # NOTE: tried extractor_args player_client=[android,ios,tv,web] to
        # dodge the GCP-IP "Sign in to confirm you're not a bot" wall —
        # verified ineffective (YouTube now blocks all InnerTube clients
        # from this datacenter IP). Reverted: the multi-client try only
        # multiplied requests against the banned IP for no benefit. A hard
        # IP ban can't be bypassed in code; it self-clears in hours, and
        # urgent videos go via manual transcript paste.
    }
    # Burner-account cookies are the one automated way past the bot wall on
    # a datacenter IP. Used only when the file exists → no-op otherwise.
    # (os/config are lazy-imported here — neither is a module-level import.)
    try:
        import os as _osc
        from .. import config as _cfg
        _ck = getattr(_cfg, "YT_COOKIES_FILE", "") or ""
        if _ck and _osc.path.exists(_ck):
            opts["cookiefile"] = _ck
    except Exception:
        pass

    _err = {"msg": ""}

    def _extract() -> dict | None:
        try:
            with YoutubeDL(opts) as ydl:
                return ydl.extract_info(url, download=False)
        except Exception as e:
            _err["msg"] = str(e)
            log.info("yt-dlp extract_info failed: %s", e)
            return None

    # Health tracking measures ONE thing: can yt-dlp still pull data
    # from YouTube at all (i.e. is the extractor / Deno-EJS path alive)?
    # That is decided by extract_info ALONE — NOT by whether this
    # particular video happens to carry ko/en captions. A video with no
    # subtitles is a normal, expected miss; counting it as a yt-dlp
    # "failure" used to make the 24h failure-rate hit 100% from a
    # handful of caption-less clips and fire a false "yt-dlp 작동 이상"
    # alert even while yt-dlp was perfectly healthy. So: extract_info
    # failure → unhealthy; extract_info success → healthy, regardless of
    # caption availability.
    from ..store import yt_dlp_health

    info = await asyncio.to_thread(_extract)
    if not info:
        # Real yt-dlp / extractor / runtime failure.
        yt_dlp_health.record_attempt(success=False)
        # Cookies configured but STILL bot-walled → they're expired/banned.
        # Flag it so the hourly health check fires a "refresh cookies" alert.
        if opts.get("cookiefile"):
            em = (_err["msg"] or "").lower()
            if "not a bot" in em or "sign in to confirm" in em:
                yt_dlp_health.record_cookie_botwall()
        return ""
    # extract_info worked → yt-dlp itself is healthy. Record success up
    # front; caption availability below doesn't change this verdict.
    yt_dlp_health.record_attempt(success=True)

    # Prefer manual subs (more accurate); fall back to auto.
    manual = info.get("subtitles") or {}
    auto = info.get("automatic_captions") or {}
    for source in (manual, auto):
        for lang in ("ko", "en"):
            tracks = source.get(lang) or []
            # Prefer json3 (clean text segments); vtt is the fallback.
            sub_url = None
            for t in tracks:
                if t.get("ext") == "json3":
                    sub_url = t.get("url"); break
            if not sub_url:
                for t in tracks:
                    if t.get("ext") == "vtt":
                        sub_url = t.get("url"); break
            if not sub_url:
                continue
            try:
                async with httpx.AsyncClient(timeout=30, follow_redirects=True,
                                             headers={"User-Agent": "Mozilla/5.0"}) as c:
                    r = await c.get(sub_url)
                    r.raise_for_status()
                    content = r.text
            except Exception as e:
                log.info("yt-dlp subtitle fetch %s/%s failed: %s",
                         "manual" if source is manual else "auto", lang, e)
                continue
            text = _parse_subtitle_payload(content)
            if text.strip():
                log.info("yt-dlp: %s/%s captions OK (%d chars)",
                         "manual" if source is manual else "auto", lang, len(text))
                # Health already recorded as success after extract_info;
                # captions are a bonus, not a separate health signal.
                return text.strip()

    # extract_info succeeded but this video has no usable ko/en captions.
    # That's a normal miss, NOT a yt-dlp failure — health was already
    # recorded as success above. Return empty so the caller routes the
    # URL to /failed (e528340) instead of saving a stub.
    log.info("yt-dlp: no usable ko/en captions for %s (yt-dlp healthy)",
             video_id)
    return ""


def _parse_subtitle_payload(content: str) -> str:
    """Extract plain text from json3 (preferred) or vtt subtitle
    payloads. json3 events look like {events: [{segs: [{utf8: '...'}]}]}.
    vtt lines have HH:MM:SS --> HH:MM:SS timestamps separated by
    blank lines from cue text."""
    import json
    # json3: real YouTube payloads start with {"wireMagic":"pb3","pens":...
    # and "events" lives ~200+ chars in, past any naive prefix check.
    # Detect by leading '{' and try a full parse — falls through to vtt
    # if parsing fails or there's no events array.
    stripped = content.lstrip()
    if stripped.startswith("{"):
        try:
            data = json.loads(stripped)
        except Exception:
            data = None
        if isinstance(data, dict) and isinstance(data.get("events"), list):
            # json3 splits captions into per-word segments where the
            # leading space lives inside the next seg's utf8 ("월간",
            # " 안테나", " 딱"). Stripping each seg destroys those
            # spaces and yields compressed garbage ("월간안테나딱"),
            # so we only strip the joined line.
            lines: list[str] = []
            for event in data["events"]:
                buf = "".join((seg.get("utf8") or "")
                              for seg in (event.get("segs") or []))
                buf = buf.strip()
                if buf:
                    lines.append(buf)
            return "\n".join(lines)
    # vtt fallback — strip timestamps and metadata
    out: list[str] = []
    for raw in content.split("\n"):
        line = raw.strip()
        if not line:
            continue
        if line.startswith("WEBVTT") or line.startswith("Kind:") or line.startswith("Language:"):
            continue
        if "-->" in line:
            continue
        if line.isdigit():  # cue number
            continue
        out.append(line)
    return "\n".join(out)


# YouTube audio-transcription fallback: cap duration so a 3h livestream
# can't run up cost/time, and so the re-encoded mp3 stays under Gemini's
# ~20MB inline limit (60min @ 32kbps mono 16kHz ≈ 14MB).
_YT_AUDIO_MAX_MIN = 60
_YT_AUDIO_MAX_BYTES = 19_000_000


async def _fetch_youtube_audio(video_id: str) -> tuple[bytes, str]:
    """Download a caption-less video's audio as low-bitrate mono mp3 so it
    fits Gemini's inline limit, for the transcription fallback. Returns
    (bytes, mime) or (b'', '') on cap-exceeded / no-ffmpeg / any failure —
    so load_youtube degrades to its existing empty-body path. Reuses the
    burner cookies if present (same bot-wall workaround as subtitles)."""
    try:
        from yt_dlp import YoutubeDL
    except Exception:
        return b"", ""
    import os as _os
    import glob as _glob
    import shutil as _shutil
    import tempfile as _tempfile
    url = f"https://www.youtube.com/watch?v={video_id}"

    def _dl() -> tuple[bytes, str, str]:
        tmpd = _tempfile.mkdtemp(prefix="ytaud_")
        opts = {
            "quiet": True, "no_warnings": True,
            "format": "bestaudio/best",
            "outtmpl": _os.path.join(tmpd, "%(id)s.%(ext)s"),
            "postprocessors": [{
                "key": "FFmpegExtractAudio",
                "preferredcodec": "mp3",
                "preferredquality": "32",   # 32 kbps
            }],
            "postprocessor_args": ["-ac", "1", "-ar", "16000"],  # mono 16kHz
        }
        try:
            from .. import config as _cfg
            _ck = getattr(_cfg, "YT_COOKIES_FILE", "") or ""
            if _ck and _os.path.exists(_ck):
                opts["cookiefile"] = _ck
        except Exception:
            pass
        try:
            with YoutubeDL(opts) as ydl:
                info = ydl.extract_info(url, download=False)
                dur = (info or {}).get("duration") or 0
                if dur and dur > _YT_AUDIO_MAX_MIN * 60:
                    return b"", "", f"too_long({int(dur)}s)"
                ydl.download([url])
            mp3s = _glob.glob(_os.path.join(tmpd, "*.mp3"))
            if not mp3s:
                return b"", "", "no_output"
            with open(mp3s[0], "rb") as f:
                data = f.read()
            if len(data) > _YT_AUDIO_MAX_BYTES:
                return b"", "", f"too_big({len(data)}B)"
            return data, "audio/mpeg", ""
        except Exception as e:
            return b"", "", str(e)[:160]
        finally:
            _shutil.rmtree(tmpd, ignore_errors=True)

    data, mime, err = await asyncio.to_thread(_dl)
    if err:
        log.info("youtube audio fetch skipped/failed (%s): %s", video_id, err)
    return data, mime


async def load_youtube(video_id: str, url: str) -> tuple[str, str, str | None]:
    title = f"YouTube {video_id}"
    description = None
    try:
        async with httpx.AsyncClient(timeout=15, follow_redirects=True,
                                     headers={"User-Agent": "Mozilla/5.0"}) as c:
            r = await c.get(f"https://www.youtube.com/oembed?url=https://www.youtube.com/watch?v={video_id}&format=json")
            if r.status_code == 200:
                title = r.json().get("title", title)
    except Exception:
        pass
    # Primary: youtube-transcript-api (fast pure-Python). Often
    # blocked on cloud IPs but free + zero deps. v1.0 of the library
    # made `list_transcripts` an instance method (was classmethod on
    # 0.x); we support both by detecting the legacy attribute first.
    try:
        if hasattr(YouTubeTranscriptApi, "list_transcripts"):
            # 0.x API
            transcripts = YouTubeTranscriptApi.list_transcripts(video_id)
            try:
                t = transcripts.find_transcript(["ko", "en"])
            except Exception:
                t = next(iter(transcripts))
            entries = t.fetch()
            text = "\n".join(e["text"] for e in entries)
        else:
            # 1.x API — instance method, returns FetchedTranscript
            # iterable of snippets with .text attribute.
            api = YouTubeTranscriptApi()
            fetched = api.fetch(video_id, languages=["ko", "en"])
            text = "\n".join(s.text for s in fetched)
        if text.strip():
            return title, text.strip(), description
    except Exception as e:
        log_msg = f"transcript_api: {e}"
    else:
        log_msg = "transcript_api: empty"

    # Fallback: yt-dlp (different API path, usually works on GCP
    # server IPs where youtube-transcript-api gets blocked). Catches
    # ~90% of remaining videos.
    yt_dlp_text = await _fetch_youtube_subs_yt_dlp(video_id)
    if yt_dlp_text:
        return title, yt_dlp_text, description

    # Final fallback: no captions anywhere (the common reason a video lands
    # here) → download the audio and let Gemini transcribe it on ANSWER_MODEL
    # (flash) so the resulting note matches normal note quality. Fully
    # guarded: any failure (cap, bot-walled download, no ffmpeg) falls back
    # to the empty-body path below. Only fires when both caption fetchers
    # missed, so it adds zero cost to the ~90% of videos that have captions.
    try:
        from .. import config as _cfg
        audio, amime = await _fetch_youtube_audio(video_id)
        if audio:
            tx = await transcribe_audio_async(
                audio, mime_type=amime, model=_cfg.ANSWER_MODEL,
                max_tokens=32768, purpose="yt_transcribe")
            if tx and tx.strip():
                log.info("youtube: Gemini audio transcription ok for %s "
                         "(%d chars)", video_id, len(tx))
                return title, tx.strip(), description
    except Exception as e:
        log.info("youtube audio transcription fallback failed for %s: %s",
                 video_id, str(e)[:160])

    # Both transcript fetchers failed. Previously we returned a stub
    # body with manual-paste instructions, but the pipeline ingested it
    # as if it were real content — so find_by_source dedup later blocked
    # any re-extraction, the stub text polluted RAG retrieval, and a
    # transient yt-dlp outage permanently corrupted every YouTube doc
    # learned in that window (the /youtube_restub_rescan cleanup was
    # built to undo exactly that). Returning empty body instead routes
    # the URL through the pipeline's existing "본문 추출 실패" path →
    # /failed with retry_payload intact, so /failed_retry replays it
    # cleanly once yt-dlp recovers (the [기술 사유] tail is preserved in
    # the docker log via _fetch_youtube_subs_yt_dlp's log.info calls).
    log.info("youtube: transcript_api + yt-dlp both failed for %s — %s",
             video_id, log_msg)
    return title, "", description


async def load_arxiv(arxiv_id: str) -> tuple[str, str, str | None]:
    """Use the arXiv Atom API: free, structured, abstract included."""
    api = f"https://export.arxiv.org/api/query?id_list={arxiv_id}"
    async with httpx.AsyncClient(timeout=20, follow_redirects=True,
                                 headers={"User-Agent": "SecondBrain"}) as c:
        r = await c.get(api)
        r.raise_for_status()
        xml = r.text
    title = _xml_field(xml, "title", skip_first=True) or f"arXiv:{arxiv_id}"
    abstract = _xml_field(xml, "summary", skip_first=True) or ""
    return title.strip()[:200], abstract.strip(), abstract.strip() or None


def _xml_field(xml: str, tag: str, skip_first: bool = False) -> str | None:
    pattern = re.compile(fr"<{tag}>(.*?)</{tag}>", re.DOTALL)
    matches = pattern.findall(xml)
    if skip_first and len(matches) > 1:
        return matches[1]
    return matches[0] if matches else None


# Auto OCR cap for sparse-text (chart-heavy) PDFs. PDFs larger than
# this trigger an inline-button prompt asking whether to OCR the rest;
# the user opts in / out per document so we never silently bill ₩200+
# on a 100-page deep-research report. Lowered 20→10 after the user
# decided most reports' essential content sits in the first ~10 pages.
import os as _os

# Vision OCR quality / cost knobs — all env-overridable for fast
# rollback if quality regression shows up on real broker reports.
SPARSE_OCR_AUTO_CAP = int(_os.getenv("OCR_AUTO_CAP", "0"))  # pages
# Image-only PDFs (PyMuPDF returns zero text) need at least a small
# Vision pass or the doc lands in /failed as empty body. 3 pages gives
# us a title + first-page summary so the doc is searchable, and the
# OCR-extend prompt picks up the rest under user control.
OCR_IMAGE_ONLY_CAP = int(_os.getenv("OCR_IMAGE_ONLY_CAP", "3"))
# Progressive OCR ([A] cost cut, 2026-05): on the AUTO sparse/image-PDF
# trigger only (not /ocr_extend), OCR the first PROBE pages, and only
# spend Vision on the remaining cap if the probe yielded enough text
# to be worth continuing. A chart-only deck where the first 3 pages
# render <300 chars of OCR text almost never has readable content on
# pages 4-7 either — saves 4 Vision calls (~₩6/doc) with low recall
# loss. Threshold 300 (was 200) leaves a margin so mixed PDFs whose
# real text only starts on page 4-5 still get continued; the rare
# false-skip is recoverable via /pending_ocr + /ocr_extend.
_OCR_PROGRESSIVE_PROBE_PAGES = int(_os.getenv("OCR_PROBE_PAGES", "3"))
_OCR_PROGRESSIVE_MIN_TEXT = int(_os.getenv("OCR_PROBE_MIN_TEXT", "300"))
OCR_DPI = int(_os.getenv("OCR_DPI", "100"))  # render DPI
OCR_SPARSE_THRESHOLD = int(_os.getenv("OCR_SPARSE_THRESHOLD", "800"))  # chars/page
# Study-notes triage: table/chart-dense pages get rendered and sent to
# Gemini Vision (which reconstructs markdown tables) instead of trusting
# get_text, which garbles tables. Capped to bound per-note Vision spend.
VISION_TABLE_CAP = int(_os.getenv("VISION_TABLE_CAP", "8"))  # pages/note


# Generic app-default titles that get baked into PDF metadata when a
# document is exported from PowerPoint / Word / Keynote / etc. without
# an explicit title set. None of these tell you what the file is about,
# so the filename ("7 Global partnering 전략.pdf") is always a better
# choice — _looks_like_title rejects these so the filename wins.
_PLACEHOLDER_TITLES = {
    # PowerPoint / Keynote / Google Slides defaults
    "powerpoint 프레젠테이션",
    "powerpoint presentation",
    "프레젠테이션",
    "presentation",
    "슬라이드 1",
    "slide 1",
    # Word / Pages defaults
    "microsoft word - document",
    "microsoft word document",
    "document",
    # Generic / iWork
    "untitled",
    "untitled document",
    "제목 없음",
    "제목없음",
    "title",
    "no title",
}


def _looks_like_title(s: str) -> bool:
    """Reject PDF metadata.title placeholders so the filename is used
    instead. Catches:
      - empty / pure-digit / pure-symbol strings
      - app-default export titles ('PowerPoint 프레젠테이션', 'Untitled', …)
      - date-only stubs: '2013년 0월 0일'
      - internal report codes: '신한투자증권20230823f', 'KB증권20240115a',
        'samsung20230101' — company name + 6+ digits + optional letter
      - bare ID tokens: '20230823rpt' (no separators, mostly digits)
    """
    s = (s or "").strip()
    if not re.search(r"[A-Za-z가-힣]{2,}", s):
        return False
    # App-default export titles — never informative, filename is better.
    if s.lower() in _PLACEHOLDER_TITLES:
        return False
    # Internal report code pattern: word + 6+ digits + optional letter,
    # no spaces. Almost always less informative than the filename.
    if re.match(r"^[가-힣A-Za-z]+\d{6,}[A-Za-z]?$", s):
        return False
    # Bare ID-ish token: no spaces, mostly digits — generic placeholder.
    if " " not in s and len(s) < 25:
        digit_ratio = sum(c.isdigit() for c in s) / max(len(s), 1)
        if digit_ratio >= 0.4:
            return False
    return True


def load_pdf(path: Path, on_stage=None) -> tuple[str, str, str | None, dict | None]:
    """Extract text from a PDF.

    Returns (title, body, hint, ocr_meta). `ocr_meta` is non-None when
    Vision OCR augmentation ran and was page-capped; bot uses it to
    offer a "OCR the remaining N pages (~₩X)" inline button.

    Tries PyMuPDF first (handles complex layouts, embedded fonts, and
    image-text overlays better) and falls back to pypdf. When the
    text-only extract looks sparse for the number of pages (chart/
    table-heavy brokerage reports etc.), augments with a Vision OCR
    pass over the first SPARSE_OCR_AUTO_CAP pages so the numbers in
    figures don't get lost."""
    title = path.stem
    body = ""
    page_count = 0
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(str(path))
        page_count = doc.page_count
        meta_title = (doc.metadata or {}).get("title")
        if meta_title and meta_title.strip() and _looks_like_title(meta_title):
            title = meta_title.strip()
        # Per-page text + structured table extraction. find_tables()
        # surfaces broker-report tables that page.get_text("text") often
        # serialises as ragged whitespace; appending the structured rows
        # gives the embedder a cleaner signal AND lets sparse-text pages
        # cross the auto-OCR threshold without spending Vision tokens.
        page_parts: list[str] = []
        for page in doc:
            t = page.get_text("text") or ""
            tables_text = _extract_pdf_tables(page)
            if tables_text:
                t = (t + "\n\n[Tables]\n" + tables_text).strip()
            page_parts.append(t)
        body = "\n\n".join(page_parts).strip()
        doc.close()
    except Exception:
        body = ""

    if not body:
        try:
            reader = PdfReader(str(path))
            page_count = page_count or len(reader.pages)
            pages = [p.extract_text() or "" for p in reader.pages]
            body = "\n\n".join(pages).strip()
            if reader.metadata and reader.metadata.title and _looks_like_title(reader.metadata.title):
                title = reader.metadata.title
        except Exception:
            pass

    ocr_meta: dict | None = None
    if not body:
        # Image-only PDF: PyMuPDF returned zero text, so without at
        # least a few OCR'd pages the doc lands in /failed as empty
        # body. OCR_IMAGE_ONLY_CAP (default 3) gives us a title +
        # first-page summary so the doc is searchable; the rest goes
        # through the user-controlled OCR-extend prompt.
        r = _ocr_pdf_pages(path, max_pages=OCR_IMAGE_ONLY_CAP,
                           on_stage=on_stage, progressive=True)
        if r["text"]:
            body = r["text"]
        if page_count > 0:
            ocr_meta = {
                "kind": "image_only",
                "applied_pages": min(page_count, OCR_IMAGE_ONLY_CAP),
                "ocrd_pages": r["ocrd"],
                "skipped_pages": r["skipped"],
                "total_pages": page_count,
                "capped": page_count > OCR_IMAGE_ONLY_CAP,
            }
    elif page_count > 0 and len(body) / max(page_count, 1) < OCR_SPARSE_THRESHOLD:
        # Sparse text/page → likely chart/table-heavy. SPARSE_OCR_AUTO_CAP
        # default 0 (2026-05): no auto OCR — every sparse PDF emits an
        # OCR-extend prompt and the user explicitly chooses OCR /
        # 텍스트만 / Cancel. Set OCR_AUTO_CAP=3 (or higher) via env to
        # restore partial auto-OCR.
        applied = min(page_count, SPARSE_OCR_AUTO_CAP)
        if SPARSE_OCR_AUTO_CAP > 0:
            r = _ocr_pdf_pages(path, max_pages=SPARSE_OCR_AUTO_CAP,
                               on_stage=on_stage, progressive=True)
            if r["text"]:
                body = body + "\n\n--- Vision OCR augmentation ---\n\n" + r["text"]
            ocrd_pages = r["ocrd"]
            skipped_pages = r["skipped"]
        else:
            ocrd_pages = 0
            skipped_pages = 0
        ocr_meta = {
            "kind": "sparse",
            "applied_pages": applied,
            "ocrd_pages": ocrd_pages,
            "skipped_pages": skipped_pages,
            "total_pages": page_count,
            "capped": page_count > SPARSE_OCR_AUTO_CAP,
        }

    return (title or path.stem)[:200], body, None, ocr_meta


def _extract_pdf_tables(page) -> str:
    """Use PyMuPDF's table finder to pull structured rows out of a
    page. Returns a TSV-ish string ready to append to the page's text,
    or '' when no usable tables are found.

    Conservative thresholds: at least 2 rows × 2 cols and ≥6 non-empty
    cells. This filters out chart-axis lines and decorative grids that
    find_tables() sometimes flags as tables, which would otherwise
    pollute the body with noise."""
    try:
        finder = page.find_tables()
    except Exception:
        return ""
    rows_all: list[str] = []
    try:
        tables = list(finder.tables) if hasattr(finder, "tables") else list(finder)
    except Exception:
        return ""
    for t in tables:
        try:
            rows = t.extract() or []
        except Exception:
            continue
        if len(rows) < 2:
            continue
        max_cols = max((len(r) for r in rows), default=0)
        if max_cols < 2:
            continue
        filled = sum(1 for r in rows for c in r if (c or "").strip())
        if filled < 6:
            continue
        rows_all.append(
            "\n".join("\t".join((c or "").strip() for c in r) for r in rows)
        )
    return "\n\n".join(rows_all)


def _page_is_blank(img_bytes: bytes, existing_text: str) -> bool:
    """Detect cover / disclaimer / blank pages so we don't waste a
    Vision call on them. Heuristic: PyMuPDF found <30 chars of text
    AND the rendered image is ≥95% near-white pixels."""
    if len(existing_text.strip()) >= 30:
        return False
    try:
        from PIL import Image
        import io as _io
        img = Image.open(_io.BytesIO(img_bytes)).convert("L")
    except Exception:
        return False
    pixels = list(img.getdata())
    if not pixels:
        return False
    # Sample to keep the scan cheap on 1275×1650-ish images.
    sample = pixels[::100]
    white = sum(1 for p in sample if p >= 240)
    return white / len(sample) >= 0.95


def _ocr_pdf_pages(path: Path, max_pages: int = 80, dpi: int = OCR_DPI,
                   start_page: int = 1, on_stage=None,
                   skip_if_text_chars: int = 1500,
                   progressive: bool = False) -> dict:
    # DPI 150 → 100: image input tokens scale roughly with pixel
    # area, so ~55% fewer input tokens per page. Vision OCR on
    # broker-report text (12-14pt body) reads identically at 100 dpi;
    # only fine print (footnotes, dense table cells) would degrade,
    # and those are usually already in the PyMuPDF text pass.
    """Render each page to PNG and ask Gemini Vision to extract text.

    Returns {text, ocrd, skipped} so callers can report actual cost
    instead of the pessimistic range size. Pages whose PyMuPDF text
    extract already exceeds `skip_if_text_chars` are skipped — the
    body already has that text via the initial PyMuPDF pass, so
    OCR-ing them again would just create duplicate chunks and bill
    Vision unnecessarily (this matters a lot on long deep-research
    papers where back-half pages are text-heavy).

    `start_page` (1-indexed inclusive) lets the resumable-OCR flow
    extend coverage beyond the auto cap by OCR-ing pages
    [start_page, start_page + max_pages - 1]."""
    try:
        import fitz  # PyMuPDF
    except Exception:
        return {"text": "", "ocrd": 0, "skipped": 0}

    try:
        doc = fitz.open(str(path))
    except Exception:
        return {"text": "", "ocrd": 0, "skipped": 0}

    # Vision OCR backend (gemini / local / hybrid) routed through
    # ocr_client. Default OCR_BACKEND=gemini keeps the original
    # inline Gemini path bit-identical; local/hybrid route to the
    # ocr-worker container via the file queue (dormant unless
    # docker-compose --profile ocr-local is up).
    from . import ocr_client as _ocr_client

    import hashlib as _hashlib
    from concurrent.futures import ThreadPoolExecutor, as_completed
    from ..store import meta as _meta
    pages_out: list[tuple[int, str]] = []  # (page_idx, text) — sorted at end
    ocrd = 0
    skipped = 0
    cached_hits = 0
    blank_skips = 0
    end_page = start_page + max_pages - 1

    # First pass (sequential, fast): decide per-page action — skip,
    # cache hit, blank, or queue for Vision. Renders happen here too
    # because PyMuPDF page objects aren't thread-safe.
    work_items: list[tuple[int, bytes, str]] = []  # (page_idx, img_bytes, img_hash)
    for i, page in enumerate(doc, 1):
        if i < start_page:
            continue
        if i > end_page:
            pages_out.append((i, f"-- Page {i}+ truncated --"))
            break
        existing_text = (page.get_text("text") or "").strip()
        if len(existing_text) >= skip_if_text_chars:
            skipped += 1
            continue
        try:
            pix = page.get_pixmap(dpi=dpi)
            img_bytes = pix.tobytes("png")
        except Exception as e:
            pages_out.append((i, f"-- Page {i} --\n[render error: {type(e).__name__}]"))
            continue
        # L: blank / boilerplate skip
        if _page_is_blank(img_bytes, existing_text):
            blank_skips += 1
            continue
        # K: cross-doc OCR cache hit — no Vision call needed
        img_hash = _hashlib.sha1(img_bytes).hexdigest()
        cached_text = _meta.ocr_cache_get(img_hash)
        if cached_text is not None:
            if cached_text.strip():
                pages_out.append((i, f"-- Page {i} --\n{cached_text}"))
            cached_hits += 1
            continue
        work_items.append((i, img_bytes, img_hash))

    # Second pass (parallel): OCR queued pages concurrently via the
    # configured backend (Gemini direct / local worker / hybrid).
    # max_workers=7 caps in-flight calls so Gemini per-min quota isn't
    # bumped even with multiple PDFs ingesting; for local mode the
    # worker itself is the bottleneck (single-process) but the queue
    # serialises naturally, so multiple threads just block on the
    # queue read — harmless.

    def _ocr_one(item):
        i, img_bytes, img_hash = item
        try:
            text = _ocr_client.ocr_one_page(img_bytes, page_idx=i)
            if text and not text.startswith("[OCR error"):
                _meta.ocr_cache_put(img_hash, text)
            return (i, text)
        except Exception as e:
            return (i, f"[OCR error: {type(e).__name__}: {e}]")

    def _run_batch(items, completed_offset, total):
        """OCR a list of work items in parallel and collect results."""
        results: list[tuple[int, str]] = []
        if not items:
            return results
        max_workers = min(len(items), 7)
        completed = completed_offset
        with ThreadPoolExecutor(max_workers=max_workers) as pool:
            futures = {pool.submit(_ocr_one, w): w for w in items}
            for fut in as_completed(futures):
                completed += 1
                if on_stage:
                    on_stage(f"Vision OCR {completed}/{total} 페이지")
                results.append(fut.result())
        return results

    if work_items:
        total = len(work_items)
        if on_stage:
            on_stage(f"Vision OCR 0/{total} 페이지")

        if progressive and len(work_items) > _OCR_PROGRESSIVE_PROBE_PAGES:
            # Probe: run the first N pages. Decide whether the
            # remainder is worth the Vision spend based on real text
            # yield (excludes [OCR error] markers).
            probe = work_items[:_OCR_PROGRESSIVE_PROBE_PAGES]
            rest = work_items[_OCR_PROGRESSIVE_PROBE_PAGES:]
            probe_results = _run_batch(probe, 0, total)
            probe_text_len = sum(
                len(t) for _, t in probe_results
                if t and not t.startswith("[OCR error")
            )
            if probe_text_len < _OCR_PROGRESSIVE_MIN_TEXT:
                log.info(
                    "progressive OCR: probe %d pages → %d chars "
                    "(threshold %d) — skipping remaining %d pages",
                    len(probe), probe_text_len,
                    _OCR_PROGRESSIVE_MIN_TEXT, len(rest),
                )
                batch_results = probe_results
            else:
                rest_results = _run_batch(
                    rest, len(probe_results), total,
                )
                batch_results = probe_results + rest_results
        else:
            batch_results = _run_batch(work_items, 0, total)

        for page_idx, text in batch_results:
            if text:
                pages_out.append((page_idx, f"-- Page {page_idx} --\n{text}"))
            ocrd += 1

    # Sort by page index so output order matches reading order.
    pages_out.sort(key=lambda t: t[0])
    pages_out_str = [t[1] for t in pages_out]
    doc.close()
    if cached_hits or blank_skips:
        log.info("ocr cache hits=%d, blank skips=%d", cached_hits, blank_skips)
    return {
        "text": "\n\n".join(pages_out_str).strip(),
        "ocrd": ocrd,
        "skipped": skipped,
    }


async def load_pdf_async(path: Path, on_stage=None) -> tuple[str, str, str | None, dict | None]:
    return await asyncio.to_thread(load_pdf, path, on_stage)


async def ocr_pdf_pages_async(path: Path, start_page: int, max_pages: int,
                              dpi: int = 150) -> str:
    """Async wrapper around _ocr_pdf_pages used by extend_pdf_ocr."""
    return await asyncio.to_thread(
        _ocr_pdf_pages, path, max_pages=max_pages, dpi=dpi,
        start_page=start_page,
    )


def _ocr_image(img_bytes: bytes, mime_type: str = "image/jpeg") -> str:
    """Extract text from a single image via Gemini Vision. Used for
    standalone photos (screenshots, table captures, news clippings).
    Cost is ~$0.0003 per image on gemini-2.5-flash-lite."""
    try:
        from google.genai import types
        from .. import config
        from ..store import cost as _cost
    except Exception:
        return ""
    client = config.make_genai_client()
    prompt = (
        "이 이미지에 포함된 모든 텍스트를 그대로 추출하세요. "
        "표는 마크다운 표 형식으로 변환하고, 단락/제목 구조를 유지하세요. "
        "설명이나 코멘트 없이 텍스트만 출력하세요."
    )
    try:
        resp = client.models.generate_content(
            model=config.SUMMARY_MODEL,
            contents=[
                types.Part.from_bytes(data=img_bytes, mime_type=mime_type),
                types.Part.from_text(text=prompt),
            ],
            config=types.GenerateContentConfig(
                temperature=0.0,
                max_output_tokens=2048,
            ),
        )
        _cost.record_resp(config.SUMMARY_MODEL, resp, purpose="ingest")
        return (resp.text or "").strip()
    except Exception:
        return ""


async def ocr_image_async(img_bytes: bytes, mime_type: str = "image/jpeg") -> str:
    return await asyncio.to_thread(_ocr_image, img_bytes, mime_type)


def _transcribe_audio(audio_bytes: bytes, mime_type: str = "audio/ogg",
                      model: str | None = None, max_tokens: int = 8192,
                      purpose: str = "ingest") -> str:
    """Gemini Audio STT. Used for Telegram voice notes and uploaded audio
    files. Cost ~₩50 per audio hour on gemini-2.5-flash-lite. Inline byte
    limit is ~20MB; longer recordings should be split client-side.

    model/max_tokens/purpose let callers override the default (e.g. the
    YouTube fallback transcribes on ANSWER_MODEL/flash for note-grade
    quality with a higher output cap)."""
    try:
        from google.genai import types
        from .. import config
        from ..store import cost as _cost
    except Exception:
        return ""
    client = config.make_genai_client()
    model = model or config.SUMMARY_MODEL
    prompt = (
        "이 오디오를 그대로 받아쓰기 하세요. 화자가 여럿이면 단락으로 "
        "구분하고, 들리는 언어 그대로(한국어/영어 등) 출력하세요. "
        "설명/코멘트 없이 받아쓰기 텍스트만 출력하세요."
    )
    try:
        resp = client.models.generate_content(
            model=model,
            contents=[
                types.Part.from_bytes(data=audio_bytes, mime_type=mime_type),
                types.Part.from_text(text=prompt),
            ],
            config=types.GenerateContentConfig(
                temperature=0.0,
                max_output_tokens=max_tokens,
            ),
        )
        _cost.record_resp(model, resp, purpose=purpose)
        return (resp.text or "").strip()
    except Exception:
        return ""


async def transcribe_audio_async(audio_bytes: bytes, mime_type: str = "audio/ogg",
                                 model: str | None = None, max_tokens: int = 8192,
                                 purpose: str = "ingest") -> str:
    return await asyncio.to_thread(
        _transcribe_audio, audio_bytes, mime_type, model, max_tokens, purpose)


async def _load_pdf_from_bytes(data: bytes, source_url: str) -> tuple[str, str, str | None]:
    """When a URL fetch returns PDF bytes (brokerage shortlinks like
    bbn.kiwoom.com → PDF redirect), save to a temp file and reuse the
    standard PDF extractor instead of trying to parse PDF as HTML.

    Drops the ocr_meta because the temp file is cleaned up before any
    user-confirmed extension OCR could run."""
    import tempfile
    from urllib.parse import urlparse, unquote
    parsed = urlparse(source_url)
    fname = unquote(Path(parsed.path).name) or "document"
    if not fname.lower().endswith(".pdf"):
        fname = (fname or "document") + ".pdf"
    tmpdir = Path(tempfile.mkdtemp(prefix="urlpdf_"))
    tmp_path = tmpdir / fname
    try:
        tmp_path.write_bytes(data)
        title, body, hint, _ocr_meta = await load_pdf_async(tmp_path)
        return title, body, hint
    finally:
        try:
            tmp_path.unlink(missing_ok=True)
            tmpdir.rmdir()
        except Exception:
            pass


def load_pptx(path: Path) -> tuple[str, str, str | None]:
    """Robust PPTX text extraction. python-pptx returns None instead
    of '' / safe defaults in several edge cases (notes_text_frame on
    slides whose notes XML is malformed, cells missing a text_frame,
    SmartArt shapes that lack a .text accessor). The user hit an
    AttributeError 'NoneType has no text' that re-queued the file on
    every restart — wrap every accessor with a None guard so a single
    odd slide can't blow up the whole load."""
    from pptx import Presentation
    prs = Presentation(str(path))
    parts: list[str] = []
    title_guess = path.stem
    for i, slide in enumerate(prs.slides, 1):
        slide_lines: list[str] = []
        for shape in slide.shapes:
            try:
                text = getattr(shape, "text", "") or ""
                if text.strip():
                    slide_lines.append(text.strip())
            except Exception:
                pass
            if getattr(shape, "has_table", False):
                try:
                    for row in shape.table.rows:
                        cells_raw = []
                        for c in row.cells:
                            t = getattr(c, "text", None)
                            if t and t.strip():
                                cells_raw.append(t.strip())
                        if cells_raw:
                            slide_lines.append(" | ".join(cells_raw))
                except Exception:
                    pass
        if getattr(slide, "has_notes_slide", False):
            try:
                ns = slide.notes_slide
                ntf = getattr(ns, "notes_text_frame", None) if ns else None
                notes = (getattr(ntf, "text", None) or "").strip() if ntf else ""
                if notes:
                    slide_lines.append(f"[Notes] {notes}")
            except Exception:
                pass
        if i == 1 and slide_lines:
            first = slide_lines[0].splitlines()[0].strip()
            if 3 <= len(first) <= 120:
                title_guess = first
        if slide_lines:
            parts.append(f"Slide {i}:\n" + "\n".join(slide_lines))
    body = "\n\n".join(parts).strip()
    return title_guess[:200], body, None


async def load_pptx_async(path: Path) -> tuple[str, str, str | None]:
    return await asyncio.to_thread(load_pptx, path)


def load_docx(path: Path) -> tuple[str, str, str | None]:
    from docx import Document
    doc = Document(str(path))
    paragraphs: list[str] = []
    for p in doc.paragraphs:
        t = p.text.strip()
        if t:
            paragraphs.append(t)
    for tbl in doc.tables:
        for row in tbl.rows:
            cells = [c.text.strip() for c in row.cells if c.text and c.text.strip()]
            if cells:
                paragraphs.append(" | ".join(cells))
    title_guess = path.stem
    if paragraphs and 3 <= len(paragraphs[0]) <= 120:
        title_guess = paragraphs[0]
    body = "\n\n".join(paragraphs).strip()
    return title_guess[:200], body, None


async def load_docx_async(path: Path) -> tuple[str, str, str | None]:
    return await asyncio.to_thread(load_docx, path)


def load_xlsx(path: Path) -> tuple[str, str, str | None]:
    """Flatten each sheet to '| col1 | col2 | ...' rows. Skips empty
    cells and rows so a wide model with sparse data doesn't waste tokens.
    Caps rows per sheet at 1500 to keep huge ledgers under control."""
    from openpyxl import load_workbook
    wb = load_workbook(filename=str(path), read_only=True, data_only=True)
    parts: list[str] = []
    MAX_ROWS_PER_SHEET = 1500
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        sheet_lines: list[str] = []
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            if i >= MAX_ROWS_PER_SHEET:
                sheet_lines.append(f"... ({MAX_ROWS_PER_SHEET}+ rows truncated)")
                break
            cells = [
                str(c).strip() for c in row
                if c is not None and str(c).strip()
            ]
            if cells:
                sheet_lines.append(" | ".join(cells))
        if sheet_lines:
            parts.append(f"## Sheet: {sheet_name}\n" + "\n".join(sheet_lines))
    wb.close()
    body = "\n\n".join(parts).strip()
    title_guess = path.stem
    return title_guess[:200], body, None


async def load_xlsx_async(path: Path) -> tuple[str, str, str | None]:
    return await asyncio.to_thread(load_xlsx, path)
